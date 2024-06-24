import csv
import logging
import os
import sys
from nltk.tokenize import word_tokenize
from nltk import pos_tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import spacy
import re

# Create a logger instance
logger = logging.getLogger('merged_logger')
logger.setLevel(logging.INFO)

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('app.log')
file_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)
logger.addHandler(file_handler)

# Load the spaCy English model
nlp = spacy.load("en_core_web_sm")

# Create a single Presentation object
combined_presentation = Presentation()

def process_script1(filename_with_identifier, combined_presentation): # Script1: Core Process Statement
    try:
        logger.info("Processing Script 1")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        file_path = os.path.join(filename_with_identifier)

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)

        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Define position and size for oval shapes
        num_cells = 7
        oval_width = Inches(2)
        oval_height = Inches(0.8)
        text_font_size = Pt(18)
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

        # Calculate the diagonal line positions
        x_step = slide_width / num_cells
        y_step = slide_height / num_cells

        # Create oval shapes with red background for the first cells in columns A and F
        for i, cell in enumerate(rows[1]):
            if i == 3:
                continue
            left = i * x_step
            top = i * y_step

            if i == 0:
                red_text = rows[1][0].split()[0]
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)

                oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + x_step, top + y_step, oval_width, oval_height)
                oval_shape.fill.solid()
                oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)

            elif i == 5:
                red_text = rows[1][5].split()[0]
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + x_step, top + y_step, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)

                oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
                oval_shape.fill.solid()
                oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)

            elif i == 1:
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + x_step, top + y_step, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 255, 0)

            elif i == 2:
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + x_step, top + y_step, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 0, 0)

            elif i == 4:
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 255, 0)

            text_frame = shape.text_frame
            text_frame.text = cell
            text_frame.paragraphs[0].font.size = text_font_size
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            oval_text_frame = oval_shape.text_frame
            oval_text_frame.text = red_text
            oval_text_frame.paragraphs[0].font.size = text_font_size
            oval_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            oval_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        verbs = []
        for i, cell in enumerate(rows[1]):
            if i == 3:
                continue
            doc = nlp(cell)
            for token in doc:
                if token.pos_ == 'VERB':
                    verbs.append(token.text)

        verbs_chunks = [verbs[i:i+5] for i in range(0, len(verbs), 5)]

        left_margin = Inches(0.5)
        top_margin = combined_presentation.slide_height - Inches(0.5)
        box_width = Inches(1.0)
        box_height = Inches(0.5)

        for chunk in verbs_chunks:
            left = left_margin
            for word in chunk:
                textbox = slide.shapes.add_textbox(left, top_margin, box_width, box_height)
                textbox.text_frame.text = word
                textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                left += box_width
                top_margin -= box_height

        logger.info(f"Script 1 saved successfully: {os.path.join(filename_without_extension + '.pptx')}")

    except Exception as e:
        logger.error(f'An error occurred in Script 1: {str(e)}')

def process_script2(filename_with_identifier, combined_presentation):  # Script2: Non-Core Process Statement
    try:
        logger.info("Processing Script 2")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        file_path = os.path.join(filename_with_identifier)

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)
            rows = rows[2:]

        slide_layout = combined_presentation.slide_layouts[5]
        slide = combined_presentation.slides.add_slide(slide_layout)

        num_cells = len(rows[0])
        node_width = Inches(2)
        node_height = Inches(0.8)
        text_font_size = Pt(18)
        left_margin = Inches(0.5)
        top_margin = Inches(1.0)
        left = left_margin
        top = top_margin
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

        x_step = slide_width / num_cells
        y_step = slide_height / num_cells

        j = 0
        while j < len(rows):
            row = rows[j]
            for i, cell in enumerate(row):
                if not cell or i == 0 or i == 3 or i == 5 or '(' in cell:
                    continue

                if i == 1 or i == 4:
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(0, 255, 0)
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height

                if i == 2:
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(255, 0, 0)
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height
            j += 1

        verbs = []
        for i, row in enumerate(rows):
            for cell in row:
                doc = nlp(cell)
                for token in doc:
                    if token.pos_ == 'VERB':
                        verbs.append(token.text)

        verbs_chunks = [verbs[i:i+5] for i in range(0, len(verbs), 5)]

        left_margin = Inches(0.5)
        top_margin = combined_presentation.slide_height - Inches(0.5)
        box_width = Inches(1.0)
        box_height = Inches(0.5)

        for chunk in verbs_chunks:
            left = left_margin
            for word in chunk:
                textbox = slide.shapes.add_textbox(left, top_margin, box_width, box_height)
                textbox.text_frame.text = word
                textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                left += box_width
                top_margin -= box_height

        logger.info(f"Script 2 saved successfully: {os.path.join(filename_without_extension + '.pptx')}")

    except Exception as e:
        logger.error(f'An error occurred in Script 2: {str(e)}')

def process_script3(filename_with_identifier, combined_presentation):  # Script3: Corporate Policy
    try:
        logger.info("Processing Script 3")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        file_path = os.path.join(filename_with_identifier)

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)
            rows = rows[1:]

        slide_layout = combined_presentation.slide_layouts[5]
        slide = combined_presentation.slides.add_slide(slide_layout)

        num_cells = len(rows[0])
        node_width = Inches(2)
        node_height = Inches(0.8)
        text_font_size = Pt(18)
        left_margin = Inches(0.5)
        top_margin = Inches(1.0)
        left = left_margin
        top = top_margin
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

        x_step = slide_width / num_cells
        y_step = slide_height / num_cells

        j = 0
        while j < len(rows):
            row = rows[j]
            for i, cell in enumerate(row):
                if not cell or i == 0 or i == 1 or i == 5 or '(' in cell:
                    continue

                if i == 2 or i == 4:
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(0, 255, 0)
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height

                if i == 3:
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(255, 0, 0)
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height
            j += 1

        verbs = []
        for i, row in enumerate(rows):
            for cell in row:
                doc = nlp(cell)
                for token in doc:
                    if token.pos_ == 'VERB':
                        verbs.append(token.text)

        verbs_chunks = [verbs[i:i+5] for i in range(0, len(verbs), 5)]

        left_margin = Inches(0.5)
        top_margin = combined_presentation.slide_height - Inches(0.5)
        box_width = Inches(1.0)
        box_height = Inches(0.5)

        for chunk in verbs_chunks:
            left = left_margin
            for word in chunk:
                textbox = slide.shapes.add_textbox(left, top_margin, box_width, box_height)
                textbox.text_frame.text = word
                textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                left += box_width
                top_margin -= box_height

        logger.info(f"Script 3 saved successfully: {os.path.join(filename_without_extension + '.pptx')}")

    except Exception as e:
        logger.error(f'An error occurred in Script 3: {str(e)}')

def process_script4(filename_with_identifier, combined_presentation):  # Script4: Business Unit Policy
    try:
        logger.info("Processing Script 4")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        file_path = os.path.join(filename_with_identifier)

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)
            rows = rows[1:]

        slide_layout = combined_presentation.slide_layouts[5]
        slide = combined_presentation.slides.add_slide(slide_layout)

        num_cells = len(rows[0])
        node_width = Inches(2)
        node_height = Inches(0.8)
        text_font_size = Pt(18)
        left_margin = Inches(0.5)
        top_margin = Inches(1.0)
        left = left_margin
        top = top_margin
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

        x_step = slide_width / num_cells
        y_step = slide_height / num_cells

        j = 0
        while j < len(rows):
            row = rows[j]
            for i, cell in enumerate(row):
                if not cell or i == 0 or i == 2 or i == 3 or '(' in cell:
                    continue

                if i == 1:
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(0, 255, 0)
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height

                if i == 4 or i == 5:
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(255, 0, 0)
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height
            j += 1

        verbs = []
        for i, row in enumerate(rows):
            for cell in row:
                doc = nlp(cell)
                for token in doc:
                    if token.pos_ == 'VERB':
                        verbs.append(token.text)

        verbs_chunks = [verbs[i:i+5] for i in range(0, len(verbs), 5)]

        left_margin = Inches(0.5)
        top_margin = combined_presentation.slide_height - Inches(0.5)
        box_width = Inches(1.0)
        box_height = Inches(0.5)

        for chunk in verbs_chunks:
            left = left_margin
            for word in chunk:
                textbox = slide.shapes.add_textbox(left, top_margin, box_width, box_height)
                textbox.text_frame.text = word
                textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                left += box_width
                top_margin -= box_height

        logger.info(f"Script 4 saved successfully: {os.path.join(filename_without_extension + '.pptx')}")

    except Exception as e:
        logger.error(f'An error occurred in Script 4: {str(e)}')

def main():
    try:
        # Check if sufficient arguments are provided
        if len(sys.argv) < 2:
            logger.error("Please provide the filename_with_identifier as a command-line argument")
            return

        filename_with_identifier = sys.argv[1]
        logger.info(f"Processing file: {filename_with_identifier}")

        process_script1(filename_with_identifier, combined_presentation)
        process_script2(filename_with_identifier, combined_presentation)
        process_script3(filename_with_identifier, combined_presentation)
        process_script4(filename_with_identifier, combined_presentation)

        output_file = os.path.join(os.path.dirname(filename_with_identifier), "Combined_Presentation.pptx")
        combined_presentation.save(output_file)
        logger.info(f"Combined Presentation saved successfully: {output_file}")

    except Exception as e:
        logger.error(f'An error occurred in the main function: {str(e)}')


    # Main function
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 RunAll.py <filename>")
        sys.exit(1)

    filename_with_identifier = sys.argv[1]

    process_script1(filename_with_identifier, combined_presentation)
    process_script2(filename_with_identifier, combined_presentation)
    process_script3(filename_with_identifier, combined_presentation)
    process_script4(filename_with_identifier, combined_presentation)
    process_script5(filename_with_identifier, combined_presentation)

    # Save the combined presentation
    combined_presentation_path = os.path.join(os.path.splitext(filename_with_identifier)[0] + '_combined.pptx')
    combined_presentation.save(combined_presentation_path)
    logger.info(f"Combined presentation saved successfully: {combined_presentation_path}")
    print(f'All scripts have been executed and their outputs have been combined into one PowerPoint presentation: "{combined_presentation_path}"')

                  
